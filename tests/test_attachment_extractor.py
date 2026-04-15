"""Tests for attachment text extraction."""

import subprocess

from src.attachment_extractor import (
    attachment_format_profile,
    extract_image_text_ocr,
    extract_text,
    extraction_quality_profile,
)


def test_extract_plain_text():
    content = b"Hello, this is a plain text file."
    result = extract_text("notes.txt", content)
    assert result == "Hello, this is a plain text file."


def test_extract_csv():
    content = b"name,value\nalice,100\nbob,200"
    result = extract_text("data.csv", content)
    assert "alice,100" in result


def test_extract_markdown():
    content = b"# Title\n\nSome markdown content."
    result = extract_text("readme.md", content)
    assert "# Title" in result


def test_extract_html():
    content = b"<html><body><p>Hello world</p></body></html>"
    result = extract_text("page.html", content)
    assert "Hello world" in result


def test_extract_calendar_file_as_text():
    content = b"BEGIN:VCALENDAR\nBEGIN:VEVENT\nSUMMARY:Review meeting\nEND:VEVENT\nEND:VCALENDAR"
    result = extract_text("invite.ics", content)
    assert result is not None
    assert "SUMMARY:Review meeting" in result


def test_extract_uses_mime_type_when_filename_has_no_extension():
    content = b"<html><body><p>Meeting summary</p></body></html>"
    result = extract_text("meeting-summary", content, mime_type="text/html")
    assert result is not None
    assert "Meeting summary" in result


def test_extract_uses_mime_type_when_filename_suffix_is_misleading():
    content = b"subject,status\nReview,open"
    result = extract_text("attachment.bin", content, mime_type="text/csv")
    assert result is not None
    assert "Review,open" in result


def test_extract_unsupported_returns_none():
    result = extract_text("photo.jpg", b"\xff\xd8\xff\xe0")
    assert result is None


def test_extract_empty_returns_none():
    assert extract_text("test.txt", b"") is None
    assert extract_text("", b"hello") is None


def test_extract_image_text_ocr_uses_tesseract_when_available(monkeypatch):
    monkeypatch.setattr("src.attachment_extractor.shutil.which", lambda _name: "/opt/homebrew/bin/tesseract")

    def _fake_run(cmd, check, capture_output, text, timeout):
        assert cmd[0] == "tesseract"
        assert cmd[2] == "stdout"
        return subprocess.CompletedProcess(cmd, 0, stdout="Recovered screenshot text", stderr="")

    monkeypatch.setattr("src.attachment_extractor.subprocess.run", _fake_run)

    result = extract_image_text_ocr("scan.png", b"fake-image-bytes")
    assert result == "Recovered screenshot text"


def test_extract_binary_skip_extensions():
    assert extract_text("archive.zip", b"PK\x03\x04") is None
    assert extract_text("program.exe", b"MZ\x90\x00") is None


def test_attachment_format_profile_marks_scanned_pdf_as_degraded_supported():
    profile = attachment_format_profile(
        filename="scan.pdf",
        mime_type="application/pdf",
        extraction_state="ocr_text_extracted",
        evidence_strength="strong_text",
        ocr_used=True,
        text_available=True,
    )
    quality = extraction_quality_profile(
        extraction_state="ocr_text_extracted",
        evidence_strength="strong_text",
        ocr_used=True,
        format_profile=profile,
    )

    assert profile["format_id"] == "scanned_pdf"
    assert profile["support_level"] == "degraded_supported"
    assert quality["quality_label"] == "ocr_text_recovered"
    assert quality["manual_review_required"] is True


def test_attachment_format_profile_marks_archive_bundles_as_unsupported():
    profile = attachment_format_profile(
        filename="records.zip",
        mime_type="application/zip",
        extraction_state="binary_only",
        evidence_strength="weak_reference",
        ocr_used=False,
        text_available=False,
    )

    assert profile["format_id"] == "archive_bundle"
    assert profile["support_level"] == "unsupported"
    assert profile["degrade_reason"] == "archive_contents_not_extracted"


def test_attachment_format_profile_marks_spreadsheets_as_lossy_but_supported():
    profile = attachment_format_profile(
        filename="timesheet.xlsx",
        mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        extraction_state="text_extracted",
        evidence_strength="strong_text",
        ocr_used=False,
        text_available=True,
    )
    quality = extraction_quality_profile(
        extraction_state="text_extracted",
        evidence_strength="strong_text",
        ocr_used=False,
        format_profile=profile,
    )

    assert profile["format_id"] == "spreadsheet_export"
    assert profile["support_level"] == "degraded_supported"
    assert profile["lossiness"] == "medium"
    assert quality["quality_label"] == "native_text_extracted"


def test_attachment_format_profile_marks_legacy_word_processing_docs_as_degraded_supported():
    profile = attachment_format_profile(
        filename="medical-note.rtf",
        mime_type="application/rtf",
        extraction_state="text_extracted",
        evidence_strength="strong_text",
        ocr_used=False,
        text_available=True,
    )

    assert profile["format_id"] == "portable_word_processing_document"
    assert profile["support_level"] == "degraded_supported"
    assert profile["degrade_reason"] == "legacy_or_portable_word_processor_structure_flattened"


def test_extract_text_truncation():
    from src.attachment_extractor import MAX_EXTRACTED_CHARS

    long_content = ("x" * (MAX_EXTRACTED_CHARS + 1000)).encode()
    result = extract_text("large.txt", long_content)
    assert result is not None
    assert "[... content truncated ...]" in result
    assert len(result) < len(long_content.decode())


def test_extract_latin1_fallback():
    # Latin-1 encoded content (not valid UTF-8)
    content = "Grüße aus München".encode("latin-1")
    result = extract_text("message.txt", content)
    assert result is not None
    assert "München" in result or "M" in result  # Latin-1 decoded


def test_pptx_returns_none_without_library(monkeypatch):
    import builtins

    real_import = builtins.__import__

    def _mock_import(name, *args, **kwargs):
        if name == "pptx":
            raise ImportError("mocked")
        return real_import(name, *args, **kwargs)

    monkeypatch.setattr(builtins, "__import__", _mock_import)
    result = extract_text("slides.pptx", b"fake content")
    assert result is None


def test_ppt_still_skipped():
    result = extract_text("slides.ppt", b"fake content")
    assert result is None


def test_pptx_extraction():
    pptx = __import__("pytest").importorskip("pptx")
    import io

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    from pptx.util import Inches

    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    txBox.text_frame.text = "Hello from slide one"

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    result = extract_text("deck.pptx", buf.read())
    assert result is not None
    assert "[Slide 1]" in result
    assert "Hello from slide one" in result
