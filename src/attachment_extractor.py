"""Extract text content from common attachment file types.

Supported formats:
- Text-based: .txt, .csv, .log, .md, .json, .xml, .yaml, .yml
- HTML: .html, .htm (converted via parse_olm._html_to_text)
- PDF: .pdf (requires PyPDF2, optional)
- Word: .docx (requires python-docx, optional)
- Excel: .xlsx (requires openpyxl, optional)
- Images: .jpg, .png, .bmp, .tiff, .webp (via ImageEmbedder, optional)

Unsupported/binary formats return None.
"""

from __future__ import annotations

import logging

from .image_embedder import _IMAGE_EXTENSIONS

logger = logging.getLogger(__name__)

MAX_EXTRACTED_CHARS = 50_000  # Truncate at ~50KB

_TEXT_EXTENSIONS = frozenset({
    ".txt", ".csv", ".log", ".md", ".json", ".xml", ".yaml", ".yml",
    ".ini", ".cfg", ".conf", ".tsv", ".rst",
})

_SKIP_EXTENSIONS = frozenset({
    ".eml", ".msg", ".zip", ".gz", ".tar", ".rar", ".7z",
    ".exe", ".dll", ".bin", ".dat", ".iso",
    ".gif", ".ico", ".svg",
    ".mp3", ".mp4", ".wav", ".avi", ".mov", ".mkv", ".flac",
    ".ppt",
})


def is_image_attachment(filename: str) -> bool:
    """Check if a filename is a supported image format for embedding."""
    if not filename:
        return False
    return _get_extension(filename) in _IMAGE_EXTENSIONS


def extract_image_embedding(filename: str, content: bytes) -> list[float] | None:
    """Encode an image attachment into a 1024-d embedding vector.

    Requires Visualized-BGE-M3 (FlagEmbedding + weight file).
    Returns None if not available or encoding fails.

    Args:
        filename: The attachment filename.
        content: The raw image bytes.

    Returns:
        1024-d embedding list, or None.
    """
    if not is_image_attachment(filename) or not content:
        return None

    try:
        from .image_embedder import ImageEmbedder

        embedder = ImageEmbedder()
        if not embedder.is_available:
            return None
        return embedder.encode_image(content)
    except Exception:  # noqa: BLE001
        logger.debug("Failed to encode image attachment: %s", filename, exc_info=True)
        return None


def extract_text(filename: str, content: bytes) -> str | None:
    """Extract readable text from an attachment.

    Args:
        filename: The attachment filename (used to determine format).
        content: The raw attachment bytes.

    Returns:
        Extracted text, or None if the format is unsupported or extraction fails.
    """
    if not filename or not content:
        return None

    ext = _get_extension(filename)

    if ext in _SKIP_EXTENSIONS or ext in _IMAGE_EXTENSIONS:
        return None

    if ext in _TEXT_EXTENSIONS:
        return _extract_plain_text(content)

    if ext in (".html", ".htm"):
        return _extract_html(content)

    if ext == ".pdf":
        return _extract_pdf(content)

    if ext == ".docx":
        return _extract_docx(content)

    if ext == ".xlsx":
        return _extract_xlsx(content)

    if ext == ".pptx":
        return _extract_pptx(content)

    return None


def _get_extension(filename: str) -> str:
    """Get lowercase file extension."""
    dot_pos = filename.rfind(".")
    if dot_pos == -1:
        return ""
    return filename[dot_pos:].lower()


def _truncate(text: str) -> str:
    """Truncate text to MAX_EXTRACTED_CHARS with a note."""
    if len(text) <= MAX_EXTRACTED_CHARS:
        return text
    return text[:MAX_EXTRACTED_CHARS] + "\n[... content truncated ...]"


def _extract_plain_text(content: bytes) -> str | None:
    """Decode text content as UTF-8 or Latin-1 fallback."""
    try:
        text = content.decode("utf-8")
    except UnicodeDecodeError:
        try:
            text = content.decode("latin-1")
        except UnicodeDecodeError:
            return None
    text = text.strip()
    return _truncate(text) if text else None


def _extract_html(content: bytes) -> str | None:
    """Convert HTML attachment to plain text."""
    try:
        html = content.decode("utf-8")
    except UnicodeDecodeError:
        try:
            html = content.decode("latin-1")
        except UnicodeDecodeError:
            return None

    from .html_converter import html_to_text as _html_to_text

    text = _html_to_text(html).strip()
    return _truncate(text) if text else None


def _extract_pdf(content: bytes) -> str | None:
    """Extract text from a PDF using PyPDF2."""
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        logger.debug("PyPDF2 not installed; skipping PDF extraction.")
        return None

    import io

    try:
        reader = PdfReader(io.BytesIO(content))
        pages = []
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                pages.append(page_text)
        text = "\n\n".join(pages).strip()
        return _truncate(text) if text else None
    except Exception:  # noqa: BLE001
        logger.debug("Failed to extract PDF text from attachment.", exc_info=True)
        return None


def _extract_docx(content: bytes) -> str | None:
    """Extract text from a Word .docx file using python-docx."""
    try:
        from docx import Document
    except ImportError:
        logger.debug("python-docx not installed; skipping DOCX extraction.")
        return None

    import io

    try:
        doc = Document(io.BytesIO(content))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n".join(paragraphs).strip()
        return _truncate(text) if text else None
    except Exception:  # noqa: BLE001
        logger.debug("Failed to extract DOCX text from attachment.", exc_info=True)
        return None


def _extract_xlsx(content: bytes) -> str | None:
    """Extract text from an Excel .xlsx file using openpyxl."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        logger.debug("openpyxl not installed; skipping XLSX extraction.")
        return None

    import io

    try:
        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        lines = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            lines.append(f"[Sheet: {sheet}]")
            for row in ws.iter_rows(values_only=True):
                cells = [str(cell) if cell is not None else "" for cell in row]
                if any(c.strip() for c in cells):
                    lines.append("\t".join(cells))
        wb.close()
        text = "\n".join(lines).strip()
        return _truncate(text) if text else None
    except Exception:  # noqa: BLE001
        logger.debug("Failed to extract XLSX text from attachment.", exc_info=True)
        return None


def _extract_pptx(content: bytes) -> str | None:
    """Extract text from a PowerPoint .pptx file using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        logger.debug("python-pptx not installed; skipping PPTX extraction.")
        return None

    import io

    try:
        prs = Presentation(io.BytesIO(content))
        lines: list[str] = []
        for slide_num, slide in enumerate(prs.slides, 1):
            lines.append(f"[Slide {slide_num}]")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            lines.append(text)
        text = "\n".join(lines).strip()
        return _truncate(text) if text else None
    except Exception:  # noqa: BLE001
        logger.debug("Failed to extract PPTX text from attachment.", exc_info=True)
        return None
